# Importing classes needed for the tool

import pandas as pd
import streamlit as st
from PIL import Image
import os
from pptx import Presentation
from io import BytesIO
import warnings
warnings.filterwarnings("ignore")
from thinkcell import Thinkcell

#Importing classes
from Analysis_sheet import Analysis_sheet
import Functions
import Charts_Functions as Charts






#Creating Object
AS = Analysis_sheet()



################# for later##########
#foot_img = Image.open('img/Periscope.png')
#####################################

#Creating sections of the website
header = st.container()
upload_section = st.container()
company_name_section = st.container()
reference_set_section = st.container()
download_section = st.container()

reference_set_list = []
reference_set_list_name = []
rev_df_list = []
sales_df_list = []
gross_margin_df_list = []
channel_df_list = []
region_df_list = []


#Title of the website
header.title("Sales Navigator Reporting Tool")

# Creates the buttons to upload the files
analysis_sheet_browse,data_loader_browse, ppt_browse = upload_section.columns(3)

analysis_sheet_loaded = analysis_sheet_browse.file_uploader("Browse Analysis Sheet",key="sales")

data_loader_loaded=data_loader_browse.file_uploader("Browse Data Loader",key="data_loader")

ppt_loaded=ppt_browse.file_uploader("Browse the ppt template",key="ppt")


if analysis_sheet_loaded is not None and data_loader_loaded is not None and ppt_loaded is not None:
    uploaded_flag=True
    #Loading the sales sheet from Analysis sheet excel as dataframe
    analysis_sheet=pd.read_excel(analysis_sheet_loaded,sheet_name="Analysis - Sales", header=5).astype(str)

    #Loading reference set sheet from Analysis sheet excel as dataframe
    analysis_sheet_reference=pd.read_excel(analysis_sheet_loaded , sheet_name="Reference Set").astype(str)

    # Loading kpi count share sheet from Analysis sheet excel as dataframe to extract referece sets name
    Referenceset_values = pd.read_excel(analysis_sheet_loaded, sheet_name="KPI Count-SHARE", header=2).columns[2:6]
    Referenceset_values_included=Functions.remove_if_is_in_list(Referenceset_values,['Reference Set 1','Reference Set 2', 'Reference Set 3',
       'Reference Set 4'])


    #Loads the ppt template name
    ppt_name = str(ppt_loaded.name)
    presentation=Presentation(ppt_loaded)
    binary_output = BytesIO()
    presentation.save(binary_output)

    #Data Loader extraction
    ## Loads data loader as dataframe
    data_loader_df = pd.read_excel(data_loader_loaded, sheet_name="DataLoader", header=34).astype(str)

    ## This function removes empty columns on dataloader
    data_loader_df = Functions.remove_values(data_loader_df, "Unnamed")

    ## Extract revenue data from dataloader

    data_loader_df_columns = data_loader_df.columns

    rev = data_loader_df.iloc[19, 2]
    rev_dict = {data_loader_df_columns[2]: float(rev)}
    rev_df = pd.DataFrame(rev_dict, index=[0])

    ## Saves revenue data as an attribute in the Analysis sheet object
    AS.rev_df = rev_df

    ## Extracts sales employee data from dataloader
    sales_employee = data_loader_df.iloc[87, 2]
    sales_employee_dict = {data_loader_df_columns[2]: float(sales_employee)}
    sales_employee_df = pd.DataFrame(sales_employee_dict, index=[0])
    ## Saves sales employee data as an attribute in the Analysis sheet object
    AS.sales_employee_df = sales_employee_df

    ## Extracts gross margin data from dataloader

    gross_margin = data_loader_df.iloc[179, 2]
    gross_margin_dict = {data_loader_df_columns[2]: float(gross_margin)}
    gross_margin_df = pd.DataFrame(gross_margin_dict, index=[0])

    ## Saves gross margin data as an attribute in the Analysis sheet object
    AS.gross_margin_df = gross_margin_df

    # Extracting Data from Analysis sheet
    ## Extracting columns names to use from analysis sheet

    kpinames = Functions.cols_extraction(analysis_sheet, 4, 6)
    company_data = Functions.cols_extraction(analysis_sheet, 139, 140)
    bu_data = Functions.cols_extraction(analysis_sheet, 140, 238)
    referenceset1 = Functions.cols_extraction(analysis_sheet, 240, 246)
    referenceset2 = Functions.cols_extraction(analysis_sheet, 246, 252)
    referenceset3 = Functions.cols_extraction(analysis_sheet, 252, 258)
    referenceset4 = Functions.cols_extraction(analysis_sheet, 258, 263)

    ## Saving KPI Names, Company Data and Bu data as an attribute of the Analysis sheet object

    AS.kpi_names = kpinames
    AS.company_data = kpinames.join(company_data)
    AS.bu_data  = kpinames.join(bu_data)

    ## Renaming the reference set columns
    AS.referenceset1 = kpinames.join(Functions.remove_end_values(referenceset1, ".4"))
    AS.referenceset2 = kpinames.join(Functions.remove_end_values(referenceset2, ".5"))
    AS.referenceset3 = kpinames.join(Functions.remove_end_values(referenceset3, ".6"))
    AS.referenceset4 = kpinames.join(Functions.remove_end_values(referenceset4, ".7"))

else:
    uploaded_flag=False
    upload_section.warning("You need to upload the Data Loader, Analysis Sheet and Sales navigator Powerpoint template")

if uploaded_flag== False:
    pass
else:
    AS.company_name=company_name_section.text_input("Input the company name")
    bu_options = Functions.bu_selection_func(AS.bu_data)
    if len(bu_options)==0:
        pass
    else:
        AS.bu_name=company_name_section.text_input("Input data cuts attributes. (eg: Company  name by region")



#Gets the name of the BU's that we have in the DataFrame
try:
    if len(bu_options)!= 0:

        bu_selection_list = company_name_section.multiselect(
            "Please select the data cuts you want to include in the final report (eg: Bu, Region, FY)", options=bu_options,
            default=list(bu_options))
        bu_selection_list.insert(0, "KPI")
        bu_selection_list.insert(0, "KPI ID")
        try:
            AS.bu_data = AS.bu_data[bu_selection_list]

        except:
            st.error("Please Select a BU")
            st.stop()
    else: pass

except:
    st.stop()



reference_dict={

    Referenceset_values[0]:AS.referenceset1,
    Referenceset_values[1]:AS.referenceset2,
    Referenceset_values[2]:AS.referenceset3,
    Referenceset_values[3]:AS.referenceset4,

     }





reference_selection_list=reference_set_section.multiselect("Please select the references set you would like to include in the ppt",
                                             Referenceset_values_included,default=Referenceset_values_included)



Refence_sets_included = dict((k, reference_dict[k]) for k in reference_selection_list if k in reference_dict)
reference_keys=list(Refence_sets_included.keys())



values = ["Average", "Worst Percentile", "Bottom Quartile", "Median", "Top Quartile", "Best Percentile"]

if len(reference_selection_list)==0:
    st.error("No reference set has been found in the analysis sheet")
elif len(reference_selection_list)==1:


    referenceset_values1 = reference_set_section.multiselect(
        "Please select the statistic measure for the "+reference_keys[0]+ " you want to use in the ppt. (eg: Bottom Quartile , Median ,Top Quartile ) ",
        values, key="reference_1", default=["Bottom Quartile", "Median", "Top Quartile"])

    referenceset_values1.insert(0, "KPI")
    referenceset_values1.insert(0, "KPI ID")


    referenceset_df1= Refence_sets_included.get(reference_keys[0])[referenceset_values1]
    reference_set_list.append(referenceset_df1)
    analysis_sheet_reference_set1 = analysis_sheet_reference.iloc[6, 1:4]



    rev_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set1, AS.rev_df, AS.company_name,
                                                           reference_keys[0]))

    analysis_sheet_reference_set_sales1 = analysis_sheet_reference.iloc[7, 1:4]
    sales_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_sales1, AS.sales_employee_df,
                                                           AS.company_name,
                                                           reference_keys[0]))

    analysis_sheet_reference_set_margin1 = analysis_sheet_reference.iloc[8, 1:4]
    gross_margin_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_margin1, AS.gross_margin_df,
                                                           AS.company_name, reference_keys[0]))

    analysis_sheet_reference_set_channel1 = analysis_sheet_reference.iloc[13:30, 0:2]
    channel_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_channel1,
                                                                                colsname=["Channel", "Value"]))

    analysis_sheet_reference_set_region1 = analysis_sheet_reference.iloc[33:50, 0:2]

    region_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_region1,
                                                                               colsname=["Region", "Value"]))




elif len(reference_selection_list)==2:

    referenceset_values1 = reference_set_section.multiselect(
        "Please select the statistic measure for the reference set 1 you want to use in the ppt. (eg: Bottom Quartile , Median ,Top Quartile ) ",
        values, key="reference_1", default=["Bottom Quartile", "Median", "Top Quartile"])

    referenceset_values1.insert(0, "KPI")
    referenceset_values1.insert(0, "KPI ID")

    referenceset_df1 = Refence_sets_included.get(reference_keys[0])[referenceset_values1]
    reference_set_list.append(referenceset_df1)
    analysis_sheet_reference_set1 = analysis_sheet_reference.iloc[6, 1:4]

    rev_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set1, AS.rev_df, AS.company_name,
                                                           reference_keys[0]))

    analysis_sheet_reference_set_sales1 = analysis_sheet_reference.iloc[7, 1:4]
    sales_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_sales1, AS.sales_employee_df,
                                                           AS.company_name,
                                                           reference_keys[0]))

    analysis_sheet_reference_set_margin1 = analysis_sheet_reference.iloc[8, 1:4]
    gross_margin_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_margin1, AS.gross_margin_df,
                                                           AS.company_name, reference_keys[0]))

    analysis_sheet_reference_set_channel1 = analysis_sheet_reference.iloc[13:30, 0:2]
    channel_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_channel1,
                                                                                colsname=["Channel", "Value"]))

    analysis_sheet_reference_set_region1 = analysis_sheet_reference.iloc[33:50, 0:2]

    region_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_region1,
                                                                               colsname=["Region", "Value"]))








    referenceset_values2 = reference_set_section.multiselect(
        "Please select the statistic measure for the reference set 2 you want to use in the ppt. (eg: Bottom Quartile , Median ,Top Quartile ) ",
        values, key="reference_2", default=["Bottom Quartile", "Median", "Top Quartile"])

    referenceset_values2.insert(0, "KPI")
    referenceset_values2.insert(0, "KPI ID")

    referenceset_df2 = Refence_sets_included.get(reference_keys[1])[referenceset_values2]
    reference2_newcols= referenceset_df2.columns
    reference2_newcols = [x +" " for x in reference2_newcols ]
    reference2_newcols[0]=reference2_newcols[0].rstrip()
    reference2_newcols[1] = reference2_newcols[1].rstrip()
    cols_dict = dict(zip(list(referenceset_df2.columns), reference2_newcols))
    referenceset_df2.rename(columns=cols_dict,inplace=True)





    reference_set_list.append(referenceset_df2)
    analysis_sheet_reference_set2 = analysis_sheet_reference.iloc[6, 8:11]
    rev_df_list.append(Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set2, AS.rev_df,
                                                                          AS.company_name, reference_keys[1]))

    analysis_sheet_reference_set_sales2 = analysis_sheet_reference.iloc[7, 8:11]
    sales_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_sales2, AS.sales_employee_df,
                                                           AS.company_name, reference_keys[1]))

    analysis_sheet_reference_set_margin2 = analysis_sheet_reference.iloc[8, 8:12]
    gross_margin_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_margin2, AS.gross_margin_df,
                                                           AS.company_name, reference_keys[1]))

    analysis_sheet_reference_set_channel2 = analysis_sheet_reference.iloc[13:30, 7:9]
    channel_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_channel2,
                                                                                colsname=["Channel", "Value"]))

    analysis_sheet_reference_set_region2 = analysis_sheet_reference.iloc[33:50, 7:9]

    region_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_region2,
                                                                               colsname=["Region", "Value"]))










elif len(reference_selection_list) == 3:

    referenceset_values1 = reference_set_section.multiselect(
        "Please select the statistic measure for the reference set 1 you want to use in the ppt. (eg: Bottom Quartile , Median ,Top Quartile ) ",
        values, key="reference_1", default=["Bottom Quartile", "Median", "Top Quartile"])

    referenceset_values1.insert(0, "KPI")
    referenceset_values1.insert(0, "KPI ID")

    referenceset_df1 = Refence_sets_included.get(reference_keys[0])[referenceset_values1]
    reference_set_list.append(referenceset_df1)
    analysis_sheet_reference_set1 = analysis_sheet_reference.iloc[6, 1:4]

    rev_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set1, AS.rev_df, AS.company_name,
                                                           reference_keys[0]))

    analysis_sheet_reference_set_sales1 = analysis_sheet_reference.iloc[7, 1:4]
    sales_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_sales1, AS.sales_employee_df,
                                                           AS.company_name,
                                                           reference_keys[0]))

    analysis_sheet_reference_set_margin1 = analysis_sheet_reference.iloc[8, 1:4]
    gross_margin_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_margin1, AS.gross_margin_df,
                                                           AS.company_name, reference_keys[0]))

    analysis_sheet_reference_set_channel1 = analysis_sheet_reference.iloc[13:30, 0:2]
    channel_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_channel1,
                                                                                colsname=["Channel", "Value"]))

    analysis_sheet_reference_set_region1 = analysis_sheet_reference.iloc[33:50, 0:2]

    region_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_region1,
                                                                               colsname=["Region", "Value"]))

    referenceset_values2 = reference_set_section.multiselect(
        "Please select the statistic measure for the reference set 2 you want to use in the ppt. (eg: Bottom Quartile , Median ,Top Quartile ) ",
        values, key="reference_2", default=["Bottom Quartile", "Median", "Top Quartile"])

    referenceset_values2.insert(0, "KPI")
    referenceset_values2.insert(0, "KPI ID")

    referenceset_df2 = Refence_sets_included.get(reference_keys[1])[referenceset_values2]
    reference2_newcols = referenceset_df2.columns
    reference2_newcols = [x + " " for x in reference2_newcols]
    reference2_newcols[0] = reference2_newcols[0].rstrip()
    reference2_newcols[1] = reference2_newcols[1].rstrip()
    cols_dict = dict(zip(list(referenceset_df2.columns), reference2_newcols))
    referenceset_df2.rename(columns=cols_dict, inplace=True)

    reference_set_list.append(referenceset_df2)
    analysis_sheet_reference_set2 = analysis_sheet_reference.iloc[6, 8:11]
    rev_df_list.append(Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set2, AS.rev_df,
                                                                          AS.company_name, reference_keys[1]))

    analysis_sheet_reference_set_sales2 = analysis_sheet_reference.iloc[7, 8:11]
    sales_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_sales2, AS.sales_employee_df,
                                                           AS.company_name, reference_keys[1]))

    analysis_sheet_reference_set_margin2 = analysis_sheet_reference.iloc[8, 8:12]
    gross_margin_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_margin2, AS.gross_margin_df,
                                                           AS.company_name, reference_keys[1]))

    analysis_sheet_reference_set_channel2 = analysis_sheet_reference.iloc[13:30, 7:9]
    channel_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_channel2,
                                                                                colsname=["Channel", "Value"]))

    analysis_sheet_reference_set_region2 = analysis_sheet_reference.iloc[33:50, 7:9]

    region_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_region2,
                                                                               colsname=["Region", "Value"]))



    referenceset_values3 = reference_set_section.multiselect(
        "Please select the statistic measure for the reference set 1 you want to use in the ppt. (eg: Bottom Quartile , Median ,Top Quartile ) ",
        values, key="reference_3", default=["Bottom Quartile", "Median", "Top Quartile"])

    referenceset_values3.insert(0, "KPI")
    referenceset_values3.insert(0, "KPI ID")

    referenceset_df3 = Refence_sets_included.get(reference_keys[2])[referenceset_values3]
    reference3_newcols = referenceset_df3.columns
    reference3_newcols = [x + "  " for x in reference3_newcols]
    reference3_newcols[0] = reference3_newcols[0].rstrip()
    reference3_newcols[1] = reference3_newcols[1].rstrip()
    cols_dict = dict(zip(list(referenceset_df3.columns), reference3_newcols))
    referenceset_df3.rename(columns=cols_dict, inplace=True)


    reference_set_list.append(referenceset_df3)

    analysis_sheet_reference_set3 = analysis_sheet_reference.iloc[6, 15:18]
    rev_df_list.append(Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set3, AS.rev_df,
                                                                          AS.company_name, reference_keys[2]))

    analysis_sheet_reference_set_sales3 = analysis_sheet_reference.iloc[7, 8:12]
    sales_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_sales3, AS.sales_employee_df,
                                                           AS.company_name, reference_keys[2]))

    analysis_sheet_reference_set_margin3 = analysis_sheet_reference.iloc[8, 8:12]
    gross_margin_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_margin3, AS.gross_margin_df,
                                                           AS.company_name, reference_keys[2]))

    analysis_sheet_reference_set_channel3 = analysis_sheet_reference.iloc[13:30, 7:9]
    channel_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_channel3,
                                                                                colsname=["Channel", "Value"]))

    analysis_sheet_reference_set_region3 = analysis_sheet_reference.iloc[33:50, 7:9]

    region_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_region3,
                                                                               colsname=["Region", "Value"]))






elif len(reference_selection_list) == 4:

    referenceset_values1 = reference_set_section.multiselect(
        "Please select the statistic measure for the reference set 1 you want to use in the ppt. (eg: Bottom Quartile , Median ,Top Quartile ) ",
        values, key="reference_1", default=["Bottom Quartile", "Median", "Top Quartile"])

    referenceset_values1.insert(0, "KPI")
    referenceset_values1.insert(0, "KPI ID")

    referenceset_df1 = Refence_sets_included.get(reference_keys[0])[referenceset_values1]
    reference_set_list.append(referenceset_df1)
    analysis_sheet_reference_set1 = analysis_sheet_reference.iloc[6, 1:4]

    rev_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set1, AS.rev_df, AS.company_name,
                                                           reference_keys[0]))

    analysis_sheet_reference_set_sales1 = analysis_sheet_reference.iloc[7, 1:4]
    sales_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_sales1, AS.sales_employee_df,
                                                           AS.company_name,
                                                           reference_keys[0]))

    analysis_sheet_reference_set_margin1 = analysis_sheet_reference.iloc[8, 1:4]
    gross_margin_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_margin1, AS.gross_margin_df,
                                                           AS.company_name, reference_keys[0]))

    analysis_sheet_reference_set_channel1 = analysis_sheet_reference.iloc[13:30, 0:2]
    channel_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_channel1,
                                                                                colsname=["Channel", "Value"]))

    analysis_sheet_reference_set_region1 = analysis_sheet_reference.iloc[33:50, 0:2]

    region_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_region1,
                                                                               colsname=["Region", "Value"]))

    referenceset_values2 = reference_set_section.multiselect(
        "Please select the statistic measure for the reference set 2 you want to use in the ppt. (eg: Bottom Quartile , Median ,Top Quartile ) ",
        values, key="reference_2", default=["Bottom Quartile", "Median", "Top Quartile"])

    referenceset_values2.insert(0, "KPI")
    referenceset_values2.insert(0, "KPI ID")

    referenceset_df2 = Refence_sets_included.get(reference_keys[1])[referenceset_values2]
    reference2_newcols = referenceset_df2.columns
    reference2_newcols = [x + " " for x in reference2_newcols]
    reference2_newcols[0] = reference2_newcols[0].rstrip()
    reference2_newcols[1] = reference2_newcols[1].rstrip()
    cols_dict = dict(zip(list(referenceset_df2.columns), reference2_newcols))
    referenceset_df2.rename(columns=cols_dict, inplace=True)


    reference_set_list.append(referenceset_df2)
    analysis_sheet_reference_set2 = analysis_sheet_reference.iloc[6, 8:11]
    rev_df_list.append(Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set2, AS.rev_df,
                                                                          AS.company_name, reference_keys[1]))

    analysis_sheet_reference_set_sales2 = analysis_sheet_reference.iloc[7, 8:11]
    sales_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_sales2, AS.sales_employee_df,
                                                           AS.company_name, reference_keys[1]))

    analysis_sheet_reference_set_margin2 = analysis_sheet_reference.iloc[8, 8:12]
    gross_margin_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_margin2, AS.gross_margin_df,
                                                           AS.company_name, reference_keys[1]))

    analysis_sheet_reference_set_channel2 = analysis_sheet_reference.iloc[13:30, 7:9]
    channel_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_channel2,
                                                                                colsname=["Channel", "Value"]))

    analysis_sheet_reference_set_region2 = analysis_sheet_reference.iloc[33:50, 7:9]

    region_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_region2,
                                                                               colsname=["Region", "Value"]))

    referenceset_values3 = reference_set_section.multiselect(
        "Please select the statistic measure for the reference set 1 you want to use in the ppt. (eg: Bottom Quartile , Median ,Top Quartile ) ",
        values, key="reference_3", default=["Bottom Quartile", "Median", "Top Quartile"])

    referenceset_values3.insert(0, "KPI")
    referenceset_values3.insert(0, "KPI ID")

    referenceset_df3 = Refence_sets_included.get(reference_keys[2])[referenceset_values3]
    reference3_newcols = referenceset_df3.columns
    reference3_newcols = [x + "  " for x in reference3_newcols]
    reference3_newcols[0] = reference3_newcols[0].rstrip()
    reference3_newcols[1] = reference3_newcols[1].rstrip()
    cols_dict = dict(zip(list(referenceset_df3.columns), reference3_newcols))
    referenceset_df3.rename(columns=cols_dict, inplace=True)

    reference_set_list.append(referenceset_df3)

    analysis_sheet_reference_set3 = analysis_sheet_reference.iloc[6, 15:18]
    rev_df_list.append(Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set3, AS.rev_df,
                                                                          AS.company_name, reference_keys[2]))

    analysis_sheet_reference_set_sales3 = analysis_sheet_reference.iloc[7, 8:12]
    sales_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_sales3, AS.sales_employee_df,
                                                           AS.company_name, reference_keys[2]))

    analysis_sheet_reference_set_margin3 = analysis_sheet_reference.iloc[8, 8:12]
    gross_margin_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_margin3, AS.gross_margin_df,
                                                           AS.company_name, reference_keys[2]))

    analysis_sheet_reference_set_channel3 = analysis_sheet_reference.iloc[13:30, 7:9]
    channel_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_channel3,
                                                                                colsname=["Channel", "Value"]))

    analysis_sheet_reference_set_region3 = analysis_sheet_reference.iloc[33:50, 7:9]

    region_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_region3,
                                                                               colsname=["Region", "Value"]))

    referenceset_values4 = reference_set_section.multiselect(
        "Please select the statistic measure for the reference set 1 you want to use in the ppt. (eg: Bottom Quartile , Median ,Top Quartile ) ",
        values, key="reference_4", default=["Bottom Quartile", "Median", "Top Quartile"])

    referenceset_values4.insert(0, "KPI")
    referenceset_values4.insert(0, "KPI ID")

    referenceset_df4 = Refence_sets_included.get(reference_keys[3])[referenceset_values3]
    reference4_newcols = referenceset_df4.columns
    reference4_newcols = [x + "   " for x in reference4_newcols]
    reference4_newcols[0] = reference4_newcols[0].rstrip()
    reference4_newcols[1] = reference4_newcols[1].rstrip()
    cols_dict = dict(zip(list(referenceset_df4.columns), reference4_newcols))
    referenceset_df4.rename(columns=cols_dict, inplace=True)

    reference_set_list.append(referenceset_df4)

    analysis_sheet_reference_set4 = analysis_sheet_reference.iloc[6, 22:25]

    rev_df_list.append(Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set4, rev_df,
                                                                          AS.company_name, reference_keys[3]))

    analysis_sheet_reference_set_sales4 = analysis_sheet_reference.iloc[7, 22:25]

    sales_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_sales4, AS.sales_employee_df,
                                                           AS.company_name, reference_keys[3]))

    analysis_sheet_reference_set_margin4 = analysis_sheet_reference.iloc[8, 22:25]
    gross_margin_df_list.append(
        Charts.Analysis_sheet_reference_set_transformation(analysis_sheet_reference_set_margin4, gross_margin_df,
                                                           AS.company_name, reference_keys[3]))

    analysis_sheet_reference_set_channel4 = analysis_sheet_reference.iloc[13:30, 21:23]
    channel_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_channel4,
                                                                                colsname=["Channel", "Value"]))

    analysis_sheet_reference_set_region4 = analysis_sheet_reference.iloc[33:50, 21:23]

    region_df_list.append(Charts.Analysis_sheet_reference_set_transformation_2(analysis_sheet_reference_set_region4,
                                                                               colsname=["Region", "Value"]))

AS.rev_df_list = rev_df_list
AS.sales_df_list = sales_df_list
AS.gross_margin_df_list = gross_margin_df_list
AS.channel_df_list = channel_df_list
AS.region_df_list = region_df_list
AS.referenceset_df_list = reference_set_list
AS.referece_set_list_names = reference_keys

try:
    emplong1 = Charts.emp_share(AS.company_data, AS.company_name, AS.bu_data, AS.referenceset1,
                                ["KPI25", "KPI27", "KPI28", "KPI29", "KPI30", "KPI31", "KPI32"])
    emplong2 = Charts.emp_share(AS.company_data, AS.company_name, AS.bu_data, AS.referenceset2,
                                ["KPI25", "KPI27", "KPI28", "KPI29", "KPI30", "KPI31", "KPI32"])
    emplong3 = Charts.emp_share(AS.company_data, AS.company_name, AS.bu_data, AS.referenceset3,
                                ["KPI25", "KPI27", "KPI28", "KPI29", "KPI30", "KPI31", "KPI32"])
    emplong4 = Charts.emp_share(AS.company_data, AS.company_name, AS.bu_data, AS.referenceset4,
                                ["KPI25", "KPI27", "KPI28", "KPI29", "KPI30", "KPI31", "KPI32"])

    empshort1 = Charts.emp_share(AS.company_data, AS.company_name, AS.bu_data, AS.referenceset1,
                                 ["KPI25", "KPI26", "KPI32"])
    empshort2 = Charts.emp_share(AS.company_data, AS.company_name, AS.bu_data, AS.referenceset2,
                                 ["KPI25", "KPI26", "KPI32"])
    empshort3 = Charts.emp_share(AS.company_data, AS.company_name, AS.bu_data, AS.referenceset3,
                                 ["KPI25", "KPI26", "KPI32"])
    empshort4 = Charts.emp_share(AS.company_data, AS.company_name, AS.bu_data, AS.referenceset4,
                                 ["KPI25", "KPI26", "KPI32"])


except:
    st.error("Please complete all the information in the form")
    st.stop()






KPI1=Charts.bar_chart_data_frame("KPI1",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI2= Charts.bar_chart_data_frame("KPI2",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI3= Charts.bar_chart_data_frame("KPI3",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI4= Charts.bar_chart_data_frame("KPI4",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI5= Charts.bar_chart_data_frame("KPI5",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI6= Charts.bar_chart_data_frame("KPI6",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI7= Charts.bar_chart_data_frame("KPI7",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI9= Charts.bar_chart_data_frame("KPI9",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI10= Charts.bar_chart_data_frame("KPI10",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI11= Charts.bar_chart_data_frame("KPI11",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI12= Charts.bar_chart_data_frame("KPI12",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI17= Charts.bar_chart_data_frame("KPI17",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI18= Charts.bar_chart_data_frame("KPI18",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI19= Charts.bar_chart_data_frame("KPI19",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI25= Charts.bar_chart_data_frame("KPI25",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI26= Charts.bar_chart_data_frame("KPI26",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI32= Charts.bar_chart_data_frame("KPI32",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)###
KPI43= Charts.bar_chart_data_frame("KPI43",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI45= Charts.bar_chart_data_frame("KPI45",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI51= Charts.bar_chart_data_frame("KPI51",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)
KPI57= Charts.bar_chart_data_frame("KPI57",AS.company_data,AS.company_name,AS.bu_data,AS.bu_name,AS.referenceset_df_list,AS.referece_set_list_names)



template_name = ppt_name

tc = Thinkcell() # create thinkcell object
tc.add_template(template_name)

# add your template
tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI1",
    dataframe=KPI1,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI2",
    dataframe=KPI2,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI3",
    dataframe=KPI3,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI4",
    dataframe=KPI4,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI5",
    dataframe=KPI5,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI6",
    dataframe=KPI6,
) # add your dataframe

# add your template
tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI7",
    dataframe=KPI7,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI9",
    dataframe=KPI9,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI10",
    dataframe=KPI10,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI11",
    dataframe=KPI11,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI12",
    dataframe=KPI12,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI17",
    dataframe=KPI17,
) # add your dataframe

# add your template
tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI18",
    dataframe=KPI18,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI19",
    dataframe=KPI19,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI25",
    dataframe=KPI25,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI26",
    dataframe=KPI26,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI32",
    dataframe=KPI32,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI43",
    dataframe=KPI43,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI45",
    dataframe=KPI45,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI51",
    dataframe=KPI51,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="KPI57",
    dataframe=KPI57,
) # add your dataframe

tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="emplong1",
    dataframe=emplong1,
)
tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="emplong2",
    dataframe=emplong2,
)
tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="emplong3",
    dataframe=emplong3,
)
tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="emplong4",
    dataframe=emplong4,
)
tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="empshort1",
    dataframe=empshort1,
)
tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="empshort2",
    dataframe=empshort2,
)
tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="empshort3",
    dataframe=empshort3,
)
tc.add_chart_from_dataframe(
    template_name=template_name,
    chart_name="empshort4",
    dataframe=empshort4,
)





for i in range(0, len(AS.rev_df_list)):

    tc.add_chart_from_dataframe(
        template_name=template_name,
        chart_name="Region"+str(i),
        dataframe=AS.region_df_list[i],
    )  # add your dataframe

    tc.add_chart_from_dataframe(
        template_name=template_name,
        chart_name="Rev" + str(i),
        dataframe=AS.rev_df_list[i],
    )  # add your dataframe

    tc.add_chart_from_dataframe(
        template_name=template_name,
        chart_name="Channel" + str(i),
        dataframe=AS.channel_df_list[i],
    )  # add your dataframe

    tc.add_chart_from_dataframe(
        template_name=template_name,
        chart_name="Sales" + str(i),
        dataframe=AS.sales_df_list[i],
    )  # add your dataframe





tc.save_ppttc("Output/template.ppttc")

tc_text=str(tc)

tc_text = tc_text.replace("None","null")
tc_text = tc_text.replace("'",'"')

tc_text = tc_text.replace("#1",'')
tc_text = tc_text.replace("#2",'')
tc_text = tc_text.replace("#3",'')
tc_text = tc_text.replace("#4",'')




flag= Functions.check_input_reference(AS.referenceset_df_list,AS.referece_set_list_names)



if  AS.bu_name =="" or AS.company_name == "" or flag== False:
    st.warning("Please fill alll fields of the form to download the ppttc file and power point presentation")

else:
    download_section.write("## Click the buttoms to download the files")
    download_section.download_button("Download pttc file", data=tc_text, file_name="template.ppttc")

    download_section.download_button(label='Download powerpoint template', data=binary_output.getvalue(),
                                     file_name=template_name)




















